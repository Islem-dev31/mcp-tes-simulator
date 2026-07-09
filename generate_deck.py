import os
import sys
import subprocess

# --- AUTO-INSTALL DEPENDENCY ---
try:
    import pptx
except ImportError:
    print("python-pptx is not installed. Installing it now...")
    try:
        subprocess.check_call([sys.executable, "-m", "pip", "install", "python-pptx"])
        import pptx
        print("python-pptx installed successfully.")
    except Exception as e:
        print(f"Error installing python-pptx: {e}")
        print("Please install it manually using: pip install python-pptx")
        sys.exit(1)

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE

# --- DESIGN & COLOR PALETTE CONFIGURATION (DARK MODERN CORPORATE STYLE) ---
# Tailored brand-aligned color system for KryoDrop (ThermaShift concept)
COLOR_DARK_BG = RGBColor(7, 19, 33)       # #071321 - Premium dark background for slides
COLOR_CARD_BG = RGBColor(16, 28, 48)      # #101C30 - Lighter card container background
COLOR_CARD_BORDER = RGBColor(27, 44, 73)  # #1B2C49 - Subtle card borders
COLOR_NAVY = RGBColor(27, 54, 93)         # #1B365D - Dark Navy (Primary brand color)
COLOR_CYAN = RGBColor(41, 182, 216)       # #29B6D8 - Neon Cyan (Brand accent color)
COLOR_WHITE = RGBColor(255, 255, 255)     # #FFFFFF - Crisp white text & headers
COLOR_MUTED_TEXT = RGBColor(156, 163, 175) # #9CA3AF - Supporting gray text
COLOR_RED = RGBColor(239, 68, 68)         # #EF4444 - High-contrast highlight for problems/risks
COLOR_GREEN = RGBColor(34, 197, 94)       # #22C55E - High-contrast highlight for viability/ROI

FONT_TITLE = "Segoe UI"
FONT_BODY = "Segoe UI"

def set_slide_background(slide, color=COLOR_DARK_BG):
    """Sets a solid background color for the slide."""
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_slide_header(slide, title, category="KRYODROP | DECISION PITCH"):
    """Adds a standard header anchor to regular slides (excludes Title/Transition slides)."""
    # Tracker / Category (Small caps style)
    cat_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.4), Inches(11.7), Inches(0.3))
    tf_cat = cat_box.text_frame
    tf_cat.word_wrap = True
    tf_cat.margin_left = tf_cat.margin_top = tf_cat.margin_right = tf_cat.margin_bottom = 0
    p_cat = tf_cat.paragraphs[0]
    p_cat.text = category.upper()
    p_cat.font.name = FONT_TITLE
    p_cat.font.size = Pt(10)
    p_cat.font.bold = True
    p_cat.font.color.rgb = COLOR_CYAN
    
    # Slide Title
    title_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.65), Inches(11.7), Inches(0.8))
    tf_title = title_box.text_frame
    tf_title.word_wrap = True
    tf_title.margin_left = tf_title.margin_top = tf_title.margin_right = tf_title.margin_bottom = 0
    p_title = tf_title.paragraphs[0]
    p_title.text = title
    p_title.font.name = FONT_TITLE
    p_title.font.size = Pt(28)
    p_title.font.bold = True
    p_title.font.color.rgb = COLOR_WHITE

def draw_card_shape(slide, left, top, width, height, bg_color=COLOR_CARD_BG, border_color=COLOR_CARD_BORDER):
    """Draws a rounded card shape with optional border, hiding default borders cleanly."""
    card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    card.fill.solid()
    card.fill.fore_color.rgb = bg_color
    
    if border_color:
        card.line.color.rgb = border_color
        card.line.width = Pt(1.5)
    else:
        card.line.fill.solid()
        card.line.fill.fore_color.rgb = bg_color
    return card

def add_content_card(slide, left, top, width, height, title, paragraphs, title_color=COLOR_CYAN, bg_color=COLOR_CARD_BG, border_color=COLOR_CARD_BORDER):
    """Creates a beautiful content card with a title and multi-paragraph body text."""
    # Draw card background
    draw_card_shape(slide, left, top, width, height, bg_color, border_color)
    
    # Overlay a text box inset inside the card boundaries
    inset = Inches(0.3)
    txBox = slide.shapes.add_textbox(left + inset, top + inset, width - (inset * 2), height - (inset * 2))
    tf = txBox.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    
    # Card Title
    p_title = tf.paragraphs[0]
    p_title.text = title
    p_title.font.name = FONT_TITLE
    p_title.font.size = Pt(16)
    p_title.font.bold = True
    p_title.font.color.rgb = title_color
    p_title.space_after = Pt(12)
    
    # Card Body Paragraphs
    for p_text in paragraphs:
        p = tf.add_paragraph()
        p.text = p_text
        p.font.name = FONT_BODY
        p.font.size = Pt(11)
        p.font.color.rgb = COLOR_WHITE
        p.space_after = Pt(8)

def add_stat_card(slide, left, top, width, height, value, label, value_color=COLOR_CYAN, bg_color=COLOR_CARD_BG):
    """Creates a statistics callout card with a large number and label below it."""
    draw_card_shape(slide, left, top, width, height, bg_color, border_color=COLOR_CARD_BORDER)
    
    inset = Inches(0.2)
    txBox = slide.shapes.add_textbox(left + inset, top + inset, width - (inset * 2), height - (inset * 2))
    tf = txBox.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    
    # Stat Value (Large bold number)
    p_val = tf.paragraphs[0]
    p_val.text = value
    p_val.font.name = FONT_TITLE
    p_val.font.size = Pt(40)
    p_val.font.bold = True
    p_val.font.color.rgb = value_color
    p_val.alignment = PP_ALIGN.CENTER
    
    # Label Description
    p_lbl = tf.add_paragraph()
    p_lbl.text = label
    p_lbl.font.name = FONT_BODY
    p_lbl.font.size = Pt(10)
    p_lbl.font.bold = True
    p_lbl.font.color.rgb = COLOR_WHITE
    p_lbl.alignment = PP_ALIGN.CENTER
    p_lbl.space_before = Pt(4)

def add_image_if_exists(slide, file_name, left, top, width, height, border_color=COLOR_CYAN):
    """Tries to find and add an image from local paths, returning True if successful."""
    paths_to_try = [
        file_name,
        os.path.join(r"c:\Users\AURES\Desktop\MCP", file_name)
    ]
    for p in paths_to_try:
        if os.path.exists(p):
            try:
                # Add thin glowing border shape behind the image
                border_rect = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left - Inches(0.05), top - Inches(0.05), width + Inches(0.1), height + Inches(0.1))
                border_rect.fill.solid()
                border_rect.fill.fore_color.rgb = COLOR_CARD_BG
                border_rect.line.color.rgb = border_color
                border_rect.line.width = Pt(1.5)
                
                # Add picture
                slide.shapes.add_picture(p, left, top, width, height)
                return True
            except Exception as e:
                print(f"Error adding image {p}: {e}")
    return False

def style_table_cell(cell, text, font_size=10, bold=False, text_color=COLOR_WHITE, bg_color=None, align=PP_ALIGN.LEFT):
    """Styles a table cell with font, color, alignment and background fill."""
    cell.text = text
    p = cell.text_frame.paragraphs[0]
    p.font.name = "Segoe UI"
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = text_color
    p.alignment = align
    if bg_color:
        cell.fill.solid()
        cell.fill.fore_color.rgb = bg_color

def create_textbox_tf(slide, left, top, width, height):
    """Helper to create a text box and return its text frame with zero margins."""
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_top = tf.margin_right = tf.margin_bottom = 0
    return tf

def build_pitch_deck():
    prs = Presentation()
    
    # Widescreen 16:9 format
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    
    blank_layout = prs.slide_layouts[6] # Blank slide layout

    # =========================================================================
    # SLIDE 1: COVER SLIDE (ThermaShift / KryoDrop)
    # =========================================================================
    slide1 = prs.slides.add_slide(blank_layout)
    set_slide_background(slide1, COLOR_DARK_BG)
    
    # Left vertical accent line
    accent_bar = slide1.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(1.0), Inches(2.0), Inches(0.12), Inches(3.5))
    accent_bar.fill.solid()
    accent_bar.fill.fore_color.rgb = COLOR_CYAN
    accent_bar.line.fill.solid()
    accent_bar.line.fill.fore_color.rgb = COLOR_CYAN
    
    # Title Text Frame
    tf1 = create_textbox_tf(slide1, Inches(1.4), Inches(1.9), Inches(7.2), Inches(4.0))
    p_main = tf1.paragraphs[0]
    p_main.text = "KryoDrop"
    p_main.font.name = FONT_TITLE
    p_main.font.size = Pt(64)
    p_main.font.bold = True
    p_main.font.color.rgb = COLOR_WHITE
    
    p_sub = tf1.add_paragraph()
    p_sub.text = "ThermaShift : Le stockage thermique intelligent pour l'agro-industrie"
    p_sub.font.name = FONT_TITLE
    p_sub.font.size = Pt(20)
    p_sub.font.bold = True
    p_sub.font.color.rgb = COLOR_CYAN
    p_sub.space_before = Pt(8)
    p_sub.space_after = Pt(20)
    
    p_desc = tf1.add_paragraph()
    p_desc.text = "Solution brevetée de stockage thermique par matériau à changement de phase (MCP-TES) pour optimiser et décarboner les chambres froides industrielles."
    p_desc.font.name = FONT_BODY
    p_desc.font.size = Pt(14)
    p_desc.font.color.rgb = COLOR_MUTED_TEXT
    
    # Logo insertion on the right if available
    # Directly insert without border since it is the transparent logo
    logo_path = os.path.join(r"c:\Users\AURES\Desktop\MCP", "logo_icon.png")
    if os.path.exists(logo_path):
        slide1.shapes.add_picture(logo_path, Inches(9.2), Inches(2.0), Inches(3.2), Inches(3.2))
    
    # Footer Details
    tf_foot = create_textbox_tf(slide1, Inches(1.4), Inches(6.1), Inches(10.0), Inches(0.5))
    p_foot = tf_foot.paragraphs[0]
    p_foot.text = "Greentech Challenge 2026  |  Pitch de Validation Industrielle"
    p_foot.font.name = FONT_BODY
    p_foot.font.size = Pt(11)
    p_foot.font.color.rgb = COLOR_CYAN

    # =========================================================================
    # SLIDE 2: THE PROBLEM (Factures Sonelgaz & Pertes)
    # =========================================================================
    slide2 = prs.slides.add_slide(blank_layout)
    set_slide_background(slide2, COLOR_DARK_BG)
    add_slide_header(slide2, "L'Explosion des Coûts & la Vulnérabilité des Stocks", "LE CONTEXTE ET LE PROBLÈME")
    
    # Column 1 (Problem 1)
    x_col1 = Inches(0.8)
    w_col = Inches(3.6)
    h_col = Inches(3.4)
    draw_card_shape(slide2, x_col1, Inches(1.6), w_col, h_col)
    # Large Number "01" at top right
    tf_num1 = create_textbox_tf(slide2, x_col1 + w_col - Inches(0.9), Inches(1.7), Inches(0.8), Inches(0.5))
    p_n1 = tf_num1.paragraphs[0]
    p_n1.text = "01"
    p_n1.font.name = FONT_TITLE
    p_n1.font.size = Pt(28)
    p_n1.font.bold = True
    p_n1.font.color.rgb = COLOR_RED
    
    tf_c1 = create_textbox_tf(slide2, x_col1 + Inches(0.25), Inches(1.8), w_col - Inches(0.5), h_col - Inches(0.5))
    p_t1 = tf_c1.paragraphs[0]
    p_t1.text = "Factures Sonelgaz HP"
    p_t1.font.name = FONT_TITLE
    p_t1.font.size = Pt(14)
    p_t1.font.bold = True
    p_t1.font.color.rgb = COLOR_CYAN
    p_t1.space_after = Pt(12)
    p_b1 = tf_c1.add_paragraph()
    p_b1.text = "Les compresseurs frigorifiques tournent à plein régime en journée, subissant les tarifs électriques Heures Pleines (HP) x3.6 plus élevés qu'en Heures Creuses (HC)."
    p_b1.font.name = FONT_BODY
    p_b1.font.size = Pt(10.5)
    p_b1.font.color.rgb = COLOR_WHITE
    p_b1.space_after = Pt(8)

    # Column 2 (Problem 2)
    x_col2 = Inches(4.8)
    draw_card_shape(slide2, x_col2, Inches(1.6), w_col, h_col)
    tf_num2 = create_textbox_tf(slide2, x_col2 + w_col - Inches(0.9), Inches(1.7), Inches(0.8), Inches(0.5))
    p_n2 = tf_num2.paragraphs[0]
    p_n2.text = "02"
    p_n2.font.name = FONT_TITLE
    p_n2.font.size = Pt(28)
    p_n2.font.bold = True
    p_n2.font.color.rgb = COLOR_RED
    
    tf_c2 = create_textbox_tf(slide2, x_col2 + Inches(0.25), Inches(1.8), w_col - Inches(0.5), h_col - Inches(0.5))
    p_t2 = tf_c2.paragraphs[0]
    p_t2.text = "Vulnérabilité Réseau"
    p_t2.font.name = FONT_TITLE
    p_t2.font.size = Pt(14)
    p_t2.font.bold = True
    p_t2.font.color.rgb = COLOR_CYAN
    p_t2.space_after = Pt(12)
    p_b2 = tf_c2.add_paragraph()
    p_b2.text = "Les coupures de courant d'été et les délestages répétitifs provoquent une hausse immédiate de la température de la chambre froide, menaçant le stock."
    p_b2.font.name = FONT_BODY
    p_b2.font.size = Pt(10.5)
    p_b2.font.color.rgb = COLOR_WHITE
    p_b2.space_after = Pt(8)

    # Column 3 (Problem 3)
    x_col3 = Inches(8.8)
    draw_card_shape(slide2, x_col3, Inches(1.6), w_col, h_col)
    tf_num3 = create_textbox_tf(slide2, x_col3 + w_col - Inches(0.9), Inches(1.7), Inches(0.8), Inches(0.5))
    p_n3 = tf_num3.paragraphs[0]
    p_n3.text = "03"
    p_n3.font.name = FONT_TITLE
    p_n3.font.size = Pt(28)
    p_n3.font.bold = True
    p_n3.font.color.rgb = COLOR_RED
    
    tf_c3 = create_textbox_tf(slide2, x_col3 + Inches(0.25), Inches(1.8), w_col - Inches(0.5), h_col - Inches(0.5))
    p_t3 = tf_c3.paragraphs[0]
    p_t3.text = "Pertes de Marchandises"
    p_t3.font.name = FONT_TITLE
    p_t3.font.size = Pt(14)
    p_t3.font.bold = True
    p_t3.font.color.rgb = COLOR_CYAN
    p_t3.space_after = Pt(12)
    p_b3 = tf_c3.add_paragraph()
    p_b3.text = "Une seule rupture de chaîne de froid entraîne la perte sèche de denrées périssables de grande valeur. Risques sanitaires et financiers colossaux pour l'exploitant."
    p_b3.font.name = FONT_BODY
    p_b3.font.size = Pt(10.5)
    p_b3.font.color.rgb = COLOR_WHITE
    p_b3.space_after = Pt(8)
    
    # Bottom Stats
    add_stat_card(slide2, left=Inches(0.8), top=Inches(5.3), width=Inches(3.6), height=Inches(1.6), value="+30 %", label="Surcoût en Heures Pleines", value_color=COLOR_RED)
    add_stat_card(slide2, left=Inches(4.8), top=Inches(5.3), width=Inches(3.6), height=Inches(1.6), value="13 Heures", label="Délestages d'été moyens supportés", value_color=COLOR_RED)
    add_stat_card(slide2, left=Inches(8.8), top=Inches(5.3), width=Inches(3.6), height=Inches(1.6), value="1.5M DA", label="Perte par incident de délestage", value_color=COLOR_RED)

    # =========================================================================
    # SLIDE 3: THE SOLUTION (Batterie PCM-TES)
    # =========================================================================
    slide3 = prs.slides.add_slide(blank_layout)
    set_slide_background(slide3, COLOR_DARK_BG)
    add_slide_header(slide3, "KryoDrop : Charger la Nuit, Libérer le Jour", "NOTRE PROPOSITION DE VALEUR")
    
    # Left Hero Box
    add_content_card(
        slide3,
        left=Inches(0.8), top=Inches(1.6), width=Inches(6.2), height=Inches(5.3),
        title="Le Cycle de Fonctionnement Intelligent",
        paragraphs=[
            "• Stockage de Froid (Nuit) : Le compresseur fonctionne de nuit pendant les Heures Creuses (tarif réduit). Le COP nocturne est supérieur de 20% grâce à la fraîcheur de l'air extérieur, permettant une solidification du MCP (paraffine) à haute efficacité énergétique.",
            "• Restitution Passive (Jour) : Le compresseur est éteint durant les Heures Pleines (8 à 12 heures). KryoDrop maintient la consigne thermique (+4°C / -18°C) par fusion contrôlée du MCP sans solliciter le compresseur.",
            "• Effacement Actif de la pointe de consommation sur le réseau électrique Sonelgaz."
        ],
        title_color=COLOR_CYAN,
        bg_color=COLOR_CARD_BG,
        border_color=COLOR_CYAN
    )
    
    # Right Stacked Cards
    # Card 1
    w_rcol = Inches(5.1)
    draw_card_shape(slide3, Inches(7.4), Inches(1.6), w_rcol, Inches(2.4))
    tf_rs1 = create_textbox_tf(slide3, Inches(7.7), Inches(1.8), w_rcol - Inches(0.6), Inches(2.0))
    p_rst1 = tf_rs1.paragraphs[0]
    p_rst1.text = "Arbitrage Énergétique Direct"
    p_rst1.font.name = FONT_TITLE
    p_rst1.font.size = Pt(15)
    p_rst1.font.bold = True
    p_rst1.font.color.rgb = COLOR_CYAN
    p_rst1.space_after = Pt(8)
    p_rsb1 = tf_rs1.add_paragraph()
    p_rsb1.text = "Déplacement de 100% de la consommation de pointe vers les heures creuses nocturnes, permettant d'économiser sur les coûts de facture variables et fixes."
    p_rsb1.font.name = FONT_BODY
    p_rsb1.font.size = Pt(10.5)
    p_rsb1.font.color.rgb = COLOR_WHITE
    
    # Card 2
    draw_card_shape(slide3, Inches(7.4), Inches(4.5), w_rcol, Inches(2.4))
    tf_rs2 = create_textbox_tf(slide3, Inches(7.7), Inches(4.7), w_rcol - Inches(0.6), Inches(2.0))
    p_rst2 = tf_rs2.paragraphs[0]
    p_rst2.text = "Sécurisation Active (13h)"
    p_rst2.font.name = FONT_TITLE
    p_rst2.font.size = Pt(15)
    p_rst2.font.bold = True
    p_rst2.font.color.rgb = COLOR_CYAN
    p_rst2.space_after = Pt(8)
    p_rsb2 = tf_rs2.add_paragraph()
    p_rsb2.text = "En cas de délestage ou de panne, la batterie thermique prend le relais automatiquement pour maintenir le froid pendant plus de 13 heures, éliminant les pertes de marchandises."
    p_rsb2.font.name = FONT_BODY
    p_rsb2.font.size = Pt(10.5)
    p_rsb2.font.color.rgb = COLOR_WHITE

    # Visual left indicators for right cards
    ind1 = slide3.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(7.4), Inches(1.6), Inches(0.08), Inches(2.4))
    ind1.fill.solid()
    ind1.fill.fore_color.rgb = COLOR_CYAN
    ind1.line.fill.solid()
    ind1.line.fill.fore_color.rgb = COLOR_CYAN
    
    ind2 = slide3.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(7.4), Inches(4.5), Inches(0.08), Inches(2.4))
    ind2.fill.solid()
    ind2.fill.fore_color.rgb = COLOR_CYAN
    ind2.line.fill.solid()
    ind2.line.fill.fore_color.rgb = COLOR_CYAN

    # =========================================================================
    # SLIDE 4: LE JUMEAU NUMÉRIQUE (Digital Twin & Algorithm)
    # =========================================================================
    slide4 = prs.slides.add_slide(blank_layout)
    set_slide_background(slide4, COLOR_DARK_BG)
    add_slide_header(slide4, "Le Jumeau Numérique de Co-Dimensionnement", "NOTRE FORCE TECHNOLOGIQUE")
    
    # Left Card: Value of Simulator
    add_content_card(
        slide4,
        left=Inches(0.8), top=Inches(1.6), width=Inches(5.2), height=Inches(5.2),
        title="Dimensionnement Sur-Mesure par Algorithme",
        paragraphs=[
            "• Modélisation Physique Rigoureuse : Notre simulateur Streamlit intègre les équations de transfert thermique transitoire à changement de phase (Stefan) et calcule la masse de MCP requise.",
            "• Base Climatique Algérienne : Intégration des données réelles des 58 wilayas d'Algérie pour évaluer les charges d'infiltration d'air lors des ouvertures de porte.",
            "• Sécurité Compresseur : Vérification de la faisabilité de la recharge nocturne par le compresseur existant.",
            "• Nomenclature instantanée (BOM) et Ratios Financiers (VAN, TRI, Payback) calculés sur-mesure pour chaque client."
        ]
    )
    
    # Right Image Card Container
    draw_card_shape(slide4, left=Inches(6.4), top=Inches(1.6), width=Inches(6.1), height=Inches(5.2), border_color=COLOR_CYAN)
    has_img4 = add_image_if_exists(slide4, "streamlit_app.png", left=Inches(6.6), top=Inches(2.1), width=Inches(5.7), height=Inches(4.2))
    if not has_img4:
        # Fallback to other screenshot
        has_img4_fb = add_image_if_exists(slide4, "Capture d'écran 2026-07-08 182424.png", left=Inches(6.6), top=Inches(2.1), width=Inches(5.7), height=Inches(4.2))
        if not has_img4_fb:
            # Placeholder text if no image
            tf_p4 = create_textbox_tf(slide4, Inches(6.6), Inches(3.0), Inches(5.7), Inches(2.0))
            p_p4 = tf_p4.paragraphs[0]
            p_p4.text = "[ Capture d'écran du Simulateur Streamlit ]\n(Simulateur d'aide à la décision technique et financière)"
            p_p4.font.name = FONT_TITLE
            p_p4.font.size = Pt(14)
            p_p4.alignment = PP_ALIGN.CENTER
            p_p4.font.color.rgb = COLOR_MUTED_TEXT

    # =========================================================================
    # SLIDE 5: INGÉNIERIE MÉCANIQUE (SolidWorks & DFM/RDM)
    # =========================================================================
    slide5 = prs.slides.add_slide(blank_layout)
    set_slide_background(slide5, COLOR_DARK_BG)
    add_slide_header(slide5, "Conception CAO & Robustesse Mécanique", "INGÉNIERIE MATÉRIELLE & DFM")
    
    # Left Card
    add_content_card(
        slide5,
        left=Inches(0.8), top=Inches(1.6), width=Inches(5.2), height=Inches(5.2),
        title="Un Design Optimisé pour la Fabrication (DFM)",
        paragraphs=[
            "• Caisson Modulaire & Système Drop-In : Conçu pour s'installer rapidement au plafond des chambres froides industrielles, sans perturber la disposition du stock.",
            "• Dégagement des Ailettes : Profil d'ailettes en étoile (8 externes, 8 internes, L=20mm, e=2mm) avec un entraxe de 50mm anti-givre pour maintenir le flux d'air.",
            "• Puits Central de 46 mm : Passage libre sécurisé pour la buse de remplissage de paraffine.",
            "• Isolation Galvanique EPDM : Manchons isolants placés aux points de contact pour éliminer la corrosion galvanique entre l'aluminium et le châssis en acier."
        ]
    )
    
    # Right Image Card Container for SolidWorks Renders
    draw_card_shape(slide5, left=Inches(6.4), top=Inches(1.6), width=Inches(6.1), height=Inches(5.2), border_color=COLOR_CYAN)
    has_img5 = add_image_if_exists(slide5, "Capture d'écran 2026-07-08 192651.png", left=Inches(6.6), top=Inches(2.1), width=Inches(5.7), height=Inches(4.2))
    if not has_img5:
        # Fallback to other screenshot
        has_img5_fb = add_image_if_exists(slide5, "Capture d'écran 2026-07-08 182424.png", left=Inches(6.6), top=Inches(2.1), width=Inches(5.7), height=Inches(4.2))
        if not has_img5_fb:
            tf_p5 = create_textbox_tf(slide5, Inches(6.6), Inches(3.0), Inches(5.7), Inches(2.0))
            p_p5 = tf_p5.paragraphs[0]
            p_p5.text = "[ Rendus SolidWorks & CAO ]\n(Batterie thermique, caisson modulaire et suspentes)"
            p_p5.font.name = FONT_TITLE
            p_p5.font.size = Pt(14)
            p_p5.alignment = PP_ALIGN.CENTER
            p_p5.font.color.rgb = COLOR_MUTED_TEXT

    # =========================================================================
    # SLIDE 6: INTELLIGENCE ÉLECTRIQUE (PLC / IHM / Automation)
    # =========================================================================
    slide6 = prs.slides.add_slide(blank_layout)
    set_slide_background(slide6, COLOR_DARK_BG)
    add_slide_header(slide6, "Contrôle-Commande & Automatisation", "INTELLIGENCE ÉLECTRIQUE")
    
    # Left Card
    add_content_card(
        slide6,
        left=Inches(0.8), top=Inches(1.6), width=Inches(5.6), height=Inches(5.2),
        title="Le Cerveau Électrique et Régulation",
        paragraphs=[
            "• Automate Industriel (PLC) : Régulation intelligente PID des flux thermiques. Basculement automatique entre les phases de charge nocturne et de décharge diurne.",
            "• Interface IHM (Écran Tactile) : Supervision locale et affichage en temps réel de l'état de charge (SOC) de la batterie thermique et des températures MCP.",
            "• Gestion de Secours Intégrée : En cas de coupure réseau, l'automate pilote le passage instantané de la ventilation sur l'onduleur de secours (UPS) pour continuer à diffuser le froid stocké.",
            "• Protection Électrique : Coffret câblé selon les normes industrielles avec sectionneurs, relais thermiques et contacteurs de puissance."
        ]
    )
    
    # Right Card: Schematic blocks representations
    draw_card_shape(slide6, left=Inches(6.8), top=Inches(1.6), width=Inches(5.7), height=Inches(5.2))
    
    # Block 1: Sensors
    card_sens = draw_card_shape(slide6, left=Inches(7.3), top=Inches(2.0), width=Inches(4.7), height=Inches(1.15), bg_color=COLOR_CARD_BG, border_color=COLOR_CARD_BORDER)
    tf_sens = create_textbox_tf(slide6, Inches(7.5), Inches(2.1), Inches(4.3), Inches(0.9))
    p_sens_title = tf_sens.paragraphs[0]
    p_sens_title.text = "1. Capteurs Thermiques (MCP & Air)"
    p_sens_title.font.name = "Segoe UI"
    p_sens_title.font.size = Pt(12)
    p_sens_title.font.bold = True
    p_sens_title.font.color.rgb = COLOR_CYAN
    p_sens = tf_sens.add_paragraph()
    p_sens.text = "Mesure en continu des températures du MCP et de l'air pour calculer l'état de charge (SOC)."
    p_sens.font.name = "Segoe UI"
    p_sens.font.size = Pt(10)
    p_sens.font.color.rgb = COLOR_WHITE
    
    # Block 2: PLC
    card_plc = draw_card_shape(slide6, left=Inches(7.3), top=Inches(3.5), width=Inches(4.7), height=Inches(1.15), bg_color=COLOR_NAVY, border_color=COLOR_CYAN)
    tf_plc = create_textbox_tf(slide6, Inches(7.5), Inches(3.6), Inches(4.3), Inches(0.9))
    p_plc_title = tf_plc.paragraphs[0]
    p_plc_title.text = "2. Automate PLC & Écran IHM"
    p_plc_title.font.name = "Segoe UI"
    p_plc_title.font.size = Pt(12)
    p_plc_title.font.bold = True
    p_plc_title.font.color.rgb = COLOR_CYAN
    p_plc = tf_plc.add_paragraph()
    p_plc.text = "Traitement des données, régulation PID, gestion des vannes et interface de contrôle opérateur."
    p_plc.font.name = "Segoe UI"
    p_plc.font.size = Pt(10)
    p_plc.font.color.rgb = COLOR_WHITE
    
    # Block 3: Actuators / UPS
    card_act = draw_card_shape(slide6, left=Inches(7.3), top=Inches(5.0), width=Inches(4.7), height=Inches(1.15), bg_color=COLOR_CARD_BG, border_color=COLOR_CARD_BORDER)
    tf_act = create_textbox_tf(slide6, Inches(7.5), Inches(5.1), Inches(4.3), Inches(0.9))
    p_act_title = tf_act.paragraphs[0]
    p_act_title.text = "3. Ventilation & Onduleur UPS"
    p_act_title.font.name = "Segoe UI"
    p_act_title.font.size = Pt(12)
    p_act_title.font.bold = True
    p_act_title.font.color.rgb = COLOR_CYAN
    p_act = tf_act.add_paragraph()
    p_act.text = "Activation de la soufflerie (60W-120W) secourue par UPS (LiFePO4 1.6 kWh) pendant 13 heures de délestage."
    p_act.font.name = "Segoe UI"
    p_act.font.size = Pt(10)
    p_act.font.color.rgb = COLOR_WHITE

    # =========================================================================
    # SLIDE 7: BUSINESS MODEL & ROI (Tableau de Rentabilité)
    # =========================================================================
    slide7 = prs.slides.add_slide(blank_layout)
    set_slide_background(slide7, COLOR_DARK_BG)
    add_slide_header(slide7, "Une Solution Rentable aux Amortissements Rapides", "BUSINESS MODEL & ROI")
    
    # Left Stats Column
    add_stat_card(slide7, left=Inches(0.8), top=Inches(1.6), width=Inches(4.0), height=Inches(1.6), value="< 3 Ans", label="Temps de Retour Simple (Payback)", value_color=COLOR_GREEN)
    add_stat_card(slide7, left=Inches(0.8), top=Inches(3.4), width=Inches(4.0), height=Inches(1.6), value="> 25 %", label="Taux de Rentabilité Interne (TRI)", value_color=COLOR_GREEN)
    add_stat_card(slide7, left=Inches(0.8), top=Inches(5.2), width=Inches(4.0), height=Inches(1.6), value="VAN > 0", label="Valeur Actuelle Nette (Projet Rentable)", value_color=COLOR_GREEN)
    
    # Right Table Card
    draw_card_shape(slide7, left=Inches(5.2), top=Inches(1.6), width=Inches(7.3), height=Inches(5.2))
    
    # Add PowerPoint Table
    rows, cols = 5, 3
    table_shape = slide7.shapes.add_table(rows, cols, Inches(5.4), Inches(2.0), Inches(6.9), Inches(4.4))
    table = table_shape.table
    
    # Set columns width
    table.columns[0].width = Inches(1.8)
    table.columns[1].width = Inches(3.4)
    table.columns[2].width = Inches(1.7)
    
    # Header Row (Navy)
    style_table_cell(table.cell(0, 0), "Poste Économique", font_size=11, bold=True, text_color=COLOR_WHITE, bg_color=COLOR_NAVY, align=PP_ALIGN.CENTER)
    style_table_cell(table.cell(0, 1), "Détails Énergétiques & Fonctionnels", font_size=11, bold=True, text_color=COLOR_WHITE, bg_color=COLOR_NAVY, align=PP_ALIGN.CENTER)
    style_table_cell(table.cell(0, 2), "Impact Budgétaire", font_size=11, bold=True, text_color=COLOR_WHITE, bg_color=COLOR_NAVY, align=PP_ALIGN.CENTER)
    
    # Row 1
    style_table_cell(table.cell(1, 0), "Arbitrage Tarifaire", font_size=10, bold=True, text_color=COLOR_WHITE, bg_color=COLOR_CARD_BG)
    style_table_cell(table.cell(1, 1), "Consommation déplacée Heures Pleines → Heures Creuses (Sonelgaz)", font_size=9.5, text_color=COLOR_MUTED_TEXT, bg_color=COLOR_CARD_BG)
    style_table_cell(table.cell(1, 2), "Économie de 30% sur le kWh", font_size=9.5, text_color=COLOR_GREEN, bg_color=COLOR_CARD_BG, align=PP_ALIGN.CENTER)
    
    # Row 2
    style_table_cell(table.cell(2, 0), "Prime de Puissance", font_size=10, bold=True, text_color=COLOR_WHITE, bg_color=COLOR_DARK_BG)
    style_table_cell(table.cell(2, 1), "Réduction de la puissance souscrite indexée sur le volume de la chambre", font_size=9.5, text_color=COLOR_MUTED_TEXT, bg_color=COLOR_DARK_BG)
    style_table_cell(table.cell(2, 2), "Gain annuel fixe", font_size=9.5, text_color=COLOR_GREEN, bg_color=COLOR_DARK_BG, align=PP_ALIGN.CENTER)
    
    # Row 3
    style_table_cell(table.cell(3, 0), "Protection Anti-Panne", font_size=10, bold=True, text_color=COLOR_WHITE, bg_color=COLOR_CARD_BG)
    style_table_cell(table.cell(3, 1), "Sécurisation active du stock de marchandises face aux coupures", font_size=9.5, text_color=COLOR_MUTED_TEXT, bg_color=COLOR_CARD_BG)
    style_table_cell(table.cell(3, 2), "Pertes de stock éliminées", font_size=9.5, text_color=COLOR_GREEN, bg_color=COLOR_CARD_BG, align=PP_ALIGN.CENTER)
    
    # Row 4
    style_table_cell(table.cell(4, 0), "Payback & Viabilité", font_size=10, bold=True, text_color=COLOR_WHITE, bg_color=COLOR_DARK_BG)
    style_table_cell(table.cell(4, 1), "Investissement amorti par les économies opérationnelles nettes", font_size=9.5, text_color=COLOR_MUTED_TEXT, bg_color=COLOR_DARK_BG)
    style_table_cell(table.cell(4, 2), "Amortissement < 3 ans", font_size=9.5, bold=True, text_color=COLOR_CYAN, bg_color=COLOR_DARK_BG, align=PP_ALIGN.CENTER)

    # =========================================================================
    # SLIDE 8: L'ÉQUIPE (ENPO-MA & ENPA Synergy)
    # =========================================================================
    slide8 = prs.slides.add_slide(blank_layout)
    set_slide_background(slide8, COLOR_DARK_BG)
    add_slide_header(slide8, "Une Équipe Pluridisciplinaire Complète", "SYNERGIE & FORCE HUMAINE")
    
    # Left Card: ENPO-MA
    add_content_card(
        slide8,
        left=Inches(0.8), top=Inches(1.6), width=Inches(5.6), height=Inches(5.2),
        title="Synergie ENPO-MA (Mécanique & Énergétique)",
        paragraphs=[
            "• Modélisation Physique & Thermique : Analyse transitoire du changement de phase solide-liquide (MCP), résolution des équations de conduction et convection forcée.",
            "• Conception Mécanique CAO : Conception complète de l'échangeur, des profilés d'ailettes en étoile SolidWorks et de la structure de supportage suspendue.",
            "• Résistance des Matériaux (RDM) : Validation de la tenue mécanique des suspentes M14, calculs de flexion (flèche médiane) des tubes de 2m, et analyses DFM de fabrication locale."
        ]
    )
    
    # Right Card: ENPA
    add_content_card(
        slide8,
        left=Inches(6.9), top=Inches(1.6), width=Inches(5.6), height=Inches(5.2),
        title="Synergie ENPA (Électrotechnique & Automatisme)",
        paragraphs=[
            "• Conception de l'Armoire Électrique : Câblage de puissance, schémas de protection et intégration des contacteurs compresseurs et de la soufflerie.",
            "• Programmation Automate (PLC) : Écriture de la logique de régulation, gestion PID des vannes, et détection automatique des coupures de courant.",
            "• Interface IHM & Supervision : Conception de l'écran tactile pour le suivi des températures, de l'état de charge (SOC) et des historiques d'exploitation.",
            "• Secours Électrique : Dimensionnement de l'UPS et du pack de batteries LiFePO4 de secours."
        ]
    )

    # --- SAVE PRESENTATION ---
    output_filename = "kryodrop_pitch_deck.pptx"
    prs.save(output_filename)
    print(f"\n[SUCCESS] Pitch deck generated successfully: '{output_filename}'")
    print(f"Absolute path: {os.path.abspath(output_filename)}")

if __name__ == "__main__":
    build_pitch_deck()
